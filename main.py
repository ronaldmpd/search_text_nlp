import os
import glob
import numpy as np
import pandas as pd
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
from fastapi import FastAPI, HTTPException, Query
from fastapi.responses import HTMLResponse
from pydantic import BaseModel
from typing import List
from sentence_transformers import SentenceTransformer, util
import re  # Asegúrate de agregar esta importación al inicio de tu archivo main.py

app = FastAPI(
    title="Buscador Semántico con BERT",
    description="API para buscar texto en documentos Word, Excel, PDF, PPTX y TXT usando BERT"
)

# --- Configuración y Carga del Modelo ---
DIRECTORIO_DOCUMENTOS = "c:\\workspace_m8_nlp\\dir_serach_docs\\"  # Asegúrate de crear esta carpeta
MODELO_BERT = "paraphrase-multilingual-MiniLM-L12-v2"

print("Cargando modelo BERT multilingüe en memoria...")
model = SentenceTransformer(MODELO_BERT)

# Almacenamiento global en memoria RAM
base_conocimiento = []
document_embeddings = None

# --- Funciones de Extracción de Texto ---
def extraer_txt(ruta):
    with open(ruta, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extraer_docx(ruta):
    doc = Document(ruta)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def extraer_pptx(ruta):
    prs = Presentation(ruta)
    texto = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texto.append(shape.text)
    return "\n".join(texto)

def extraer_pdf(ruta):
    reader = PdfReader(ruta)
    texto = []
    for page in reader.pages:
        txt = page.extract_text()
        if txt:
            texto.append(txt)
    return "\n".join(texto)

def extraer_excel(ruta):
    try:
        excel_data = pd.read_excel(ruta, sheet_name=None)
        texto = []
        for sheet_name, df in excel_data.items():
            sheet_text = df.astype(str).to_string(index=False, header=False)
            texto.append(f"Hoja {sheet_name}:\n{sheet_text}")
        return "\n".join(texto)
    except Exception:
        return ""

def fragmentar_texto(texto, max_palabras=150):
    palabras = texto.split()
    return [" ".join(palabras[i:i + max_palabras]) for i in range(0, len(palabras), max_palabras)]

# --- Modelos de Datos Pydantic ---
class StatusResponse(BaseModel):
    mensaje: str
    total_fragmentos: int

class ResultadoBusqueda(BaseModel):
    archivo: str
    tipo: str
    similitud: float
    fragmento: str

# --- Endpoints de la API ---

@app.get("/", response_class=HTMLResponse, tags=["Interfaz Web"])
def leer_index():
    """Servir la página principal del buscador HTML."""
    if not os.path.exists("index.html"):
        raise HTTPException(
            status_code=404, 
            detail="Archivo 'index.html' no encontrado en el directorio raíz."
        )
    
    with open("index.html", "r", encoding="utf-8") as f:
        return f.read()


@app.post("/indexar", response_model=StatusResponse, tags=["Indexación"])
def indexar_documentos():
    """
    Escanea el directorio configurado, extrae el texto de los archivos
    y genera los embeddings de BERT.
    """
    global base_conocimiento, document_embeddings
    
    if not os.path.exists(DIRECTORIO_DOCUMENTOS):
        os.makedirs(DIRECTORIO_DOCUMENTOS)
        print(f"Directorio creado automáticamente: '{DIRECTORIO_DOCUMENTOS}'")
    
    nuevos_fragmentos = []
    extensiones = ['*.txt', '*.docx', '*.pptx', '*.pdf', '*.xlsx', '*.xls']
    archivos = []
    
    for ext in extensiones:
        archivos.extend(glob.glob(os.path.join(DIRECTORIO_DOCUMENTOS, ext)))
        
    for archivo in archivos:
        nombre_archivo = os.path.basename(archivo)
        ext = os.path.splitext(archivo)[1].lower()
        texto_extraido = ""
        
        try:
            if ext == '.txt': texto_extraido = extraer_txt(archivo)
            elif ext == '.docx': texto_extraido = extraer_docx(archivo)
            elif ext == '.pptx': texto_extraido = extraer_pptx(archivo)
            elif ext == '.pdf': texto_extraido = extraer_pdf(archivo)
            elif ext in ['.xlsx', '.xls']: texto_extraido = extraer_excel(archivo)
            
            if texto_extraido.strip():
                fragmentos = fragmentar_texto(texto_extraido)
                for frag in fragmentos:
                    nuevos_fragmentos.append({
                        "texto": frag,
                        "archivo": nombre_archivo,
                        "tipo": ext
                    })
        except Exception as e:
            print(f"Error procesando {nombre_archivo}: {e}")

    if not nuevos_fragmentos:
        raise HTTPException(
            status_code=400, 
            detail="No se encontraron documentos válidos en la carpeta o están completamente vacíos."
        )
    
    # Actualizar estado global
    base_conocimiento = nuevos_fragmentos
    textos_indexar = [item['texto'] for item in base_conocimiento]
    
    # Generar embeddings
    document_embeddings = model.encode(textos_indexar, convert_to_tensor=True)
    
    return {
        "mensaje": "Indexación completada con éxito.",
        "total_fragmentos": len(base_conocimiento)
    }


@app.get("/buscar", response_model=List[ResultadoBusqueda], tags=["Búsqueda"])
def buscar_texto(
    q: str = Query(..., description="Texto, nombre completo o número de carnet a buscar"), 
    top_k: int = Query(3, description="Número de resultados a devolver")
):
    """
    Realiza una búsqueda híbrida y recorta el fragmento para mostrar únicamente
    las 2 líneas más relevantes donde aparece el término buscado.
    """
    global base_conocimiento, document_embeddings
    
    if not base_conocimiento or document_embeddings is None:
        raise HTTPException(
            status_code=400, 
            detail="La base de conocimiento está vacía. Indexa primero los documentos."
        )
    
    query_limpia = q.strip().lower()
    es_posible_carnet = re.search(r'\b\d{4,10}\b', query_limpia)
    
    # Esta función auxiliar se encargará de extraer solo las 2 líneas del match
    def recortar_a_dos_lineas(texto_completo, termino_buscado):
        # Separamos el fragmento en líneas reales
        lineas = [linea.strip() for linea in texto_completo.split('\n') if linea.strip()]
        
        for i, linea in enumerate(lineas):
            if termino_buscado in linea.lower():
                # Si es la última línea, tomamos la anterior y la actual
                if i == len(lineas) - 1 and i > 0:
                    return f"{lineas[i-1]}\n{lineas[i]}"
                # Por defecto, tomamos la línea del match y la que le sigue (2 líneas)
                elif i < len(lineas) - 1:
                    return f"{lineas[i]}\n{lineas[i+1]}"
                # Si el documento solo tiene una línea
                return linea
        
        # Salvaguarda: si no encuentra la línea exacta por separación de caracteres,
        # devuelve los primeros 200 caracteres del fragmento
        return texto_completo[:200]

    # 1. PASO HÍBRIDO (Match exacto por Carnet o Nombre)    
    resultados_directos = []
    termino_exacto = es_posible_carnet.group(0) if es_posible_carnet else query_limpia

    for idx, item in enumerate(base_conocimiento):
        texto_fragmento = item["texto"].lower()
        
        if es_posible_carnet and termino_exacto in texto_fragmento:
            resultados_directos.append((idx, 1.0))
        elif len(query_limpia) > 8 and query_limpia in texto_fragmento:
            resultados_directos.append((idx, 0.95))

    if resultados_directos:
        resultados_directos = sorted(resultados_directos, key=lambda x: x[1], reverse=True)[:top_k]
        resultados = []
        for idx, score in resultados_directos:
            fragmento_recortado = recortar_a_dos_lineas(base_conocimiento[idx]["texto"], termino_exacto)
            resultados.append({
                "archivo": base_conocimiento[idx]["archivo"],
                "tipo": base_conocimiento[idx]["tipo"],
                "similitul": score,  # Mantiene compatibilidad con tu HTML anterior
                "similitud": score,
                "fragmento": fragmento_recortado
            })
        return resultados

    # 2. PASO SEMÁNTICO (BERT - Si no hay coincidencia exacta de palabras)    
    k = min(top_k, len(base_conocimiento))
    query_embedding = model.encode(q, convert_to_tensor=True)
    cos_scores = util.cos_sim(query_embedding, document_embeddings)[0]
    
    top_results = np.argpartition(-cos_scores.cpu(), k)[:k]
    top_results = top_results[np.argsort(-cos_scores[top_results].cpu())]
    
    resultados = []
    for idx in top_results:
        # Para BERT, intentamos buscar alguna de las palabras de la consulta en las líneas
        # tomamos la primera palabra significativa (mayor a 3 letras) para el recorte visual
        palabras = [p for p in query_limpia.split() if len(p) > 3]
        palabra_clave = palabras[0] if palabras else query_limpia
        
        #fragmento_recortado = recortar_a_dos_lineas(base_conocimiento[idx]["texto"], palabra_clave)
        fragmento_recortado = base_conocimiento[idx]["texto"]

        resultados.append({
            "archivo": base_conocimiento[idx]["archivo"],
            "tipo": base_conocimiento[idx]["tipo"],
            "similitul": round(float(cos_scores[idx].item()), 4),
            "similitud": round(float(cos_scores[idx].item()), 4),
            "fragmento": fragmento_recortado
        })
        
    return resultados

# Evento automático para indexar documentos al arrancar el backend
@app.on_event("startup")
def auto_indexar():
    try:
        print("Intentando indexación automática inicial...")
        indexar_documentos()
    except Exception as e:
        print(f"Aviso: No se pudo auto-indexar al arrancar (Directorio vacío o sin formatos soportados): {e}")


#uvicorn main:app --reload --port 8080